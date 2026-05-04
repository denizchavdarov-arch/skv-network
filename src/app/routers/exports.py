"""Export endpoints: Markdown, PPTX, and Cube Extraction."""
from fastapi import APIRouter, Request
from fastapi.responses import PlainTextResponse, FileResponse
from app.routers.entries import cubes_library
import json
import urllib.request as _req

router = APIRouter()

POLZA_KEY = "pza_K738KdM_Cm2HYltwAvCLi3Uw9n8U5Rfo"

@router.get("/api/export/markdown/{cube_id}")
async def export_markdown(cube_id: str):
    cube = cubes_library.get(cube_id, {})
    if not cube:
        return {"error": "Cube not found"}
    c = cube.get("content", cube)
    if isinstance(c, str):
        c = json.loads(c)
    title = c.get("title", cube_id)
    typ = c.get("type", "N/A")
    prio = c.get("priority", "N/A")
    stat = c.get("status", "N/A")
    trig = ", ".join(c.get("trigger_intent", []))
    rules = c.get("rules", [])
    rat = c.get("rationale", "")
    exs = c.get("examples", [])
    md = "# " + title + "\n\n"
    md += "Type: " + typ + " | Priority: " + str(prio) + " | Status: " + stat + "\n\n"
    md += "Triggers: " + trig + "\n\n"
    md += "## Rules\n\n"
    for i, rule in enumerate(rules, 1):
        md += str(i) + ". " + rule + "\n"
    if rat:
        md += "\n## Rationale\n\n" + rat + "\n"
    if exs:
        md += "\n## Examples\n\n"
        for ex in exs:
            md += "- " + ex + "\n"
    return PlainTextResponse(md, media_type="text/markdown")

@router.get("/api/export/pptx/{cube_id}")
async def export_pptx(cube_id: str):
    cube = cubes_library.get(cube_id, {})
    if not cube:
        return {"error": "Cube not found"}
    c = cube.get("content", cube)
    if isinstance(c, str):
        c = json.loads(c)
    from pptx import Presentation
    from pptx.util import Inches, Pt
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = c.get("title", cube_id)
    slide.placeholders[1].text = "Type: " + str(c.get("type", "")) + " | Priority: " + str(c.get("priority", ""))
    if c.get("rules"):
        slide2 = prs.slides.add_slide(prs.slide_layouts[1])
        slide2.shapes.title.text = "Rules"
        txBox = slide2.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(5))
        tf = txBox.text_frame
        for rule in c["rules"]:
            p = tf.add_paragraph()
            p.text = "• " + rule
            p.font.size = Pt(18)
    fname = "/app/uploads/" + cube_id + ".pptx"
    prs.save(fname)
    return FileResponse(fname, filename=cube_id + ".pptx")

@router.post("/api/experience/extract")
async def extract_cubes(request: Request):
    """Извлекает кубики SKV из текста диалога."""
    data = await request.json()
    text = data.get("text", "")[:5000]
    if not text:
        return {"error": "Text required"}

    prompt = f"""Extract useful experience cubes from this dialogue. 
Return a JSON array of cubes following SKV format:
{{"cubes": [{{"cube_id": "cube_exp_...", "type": "experience", "priority": 3, "version": "1.0.0", "title": "...", "trigger_intent": [...], "rules": ["Rule 1", "Rule 2"], "rationale": "...", "examples": [...], "source": "Extracted from dialogue"}}]}}

Dialogue:
{text}

Extract 1-5 cubes. Focus on reusable engineering lessons, debugging tips, architecture patterns. Be specific, not generic."""

    body = json.dumps({
        "model": "deepseek/deepseek-v4-flash",
        "messages": [{"role": "user", "content": prompt}],
        "temperature": 0.3,
        "max_tokens": 2000
    }).encode()

    for attempt in range(2):
        try:
            req = _req.Request("https://api.polza.ai/v1/chat/completions", data=body, headers={
                "Content-Type": "application/json",
                "Authorization": f"Bearer {POLZA_KEY}"
            })
            resp = _req.urlopen(req, timeout=120)
            answer = json.loads(resp.read())["choices"][0]["message"]["content"]
            # Extract JSON from response
            if "```json" in answer:
                answer = answer.split("```json")[1].split("```")[0]
            elif "```" in answer:
                answer = answer.split("```")[1].split("```")[0]
            cubes_data = json.loads(answer)
            return cubes_data
        except Exception as e:
            if attempt == 1:
                return {"error": str(e)[:200], "raw_answer": answer if 'answer' in dir() else ""}
    return {"error": "Empty after 2 attempts"}

@router.get("/api/system/status")
async def system_status():
    from app.routers.entries import cubes_library, get_cubes_count
    import os
    
    # Статус экстрактора
    extractor_running = os.system("pgrep -f extract_mass.py > /dev/null") == 0
    
    # Последние строки лога
    log_tail = ""
    try:
        with open("/root/skv-core/data/extract_mass.log") as f:
            lines = f.readlines()
            log_tail = "".join(lines[-5:])
    except:
        log_tail = "Log not available"
    
    return {
        "service": "SKV Network",
        "version": "2.0",
        "cubes_total": get_cubes_count(),
        "cubes_in_memory": len(cubes_library),
        "qdrant_points": "check via /api/qdrant/count",
        "extractor": {
            "running": extractor_running,
            "last_lines": log_tail
        },
        "uptime": "since restart"
    }

@router.post("/api/tools/create-file")
async def create_file(request: Request):
    """Создаёт файл на сервере и возвращает ссылку для скачивания."""
    import os
    data = await request.json()
    filename = data.get("filename", "output.txt")
    content = data.get("content", "")
    filetype = data.get("type", "txt")
    
    if not content:
        return {"error": "Content required"}
    
    # Определяем расширение
    ext_map = {"html": "html", "txt": "txt", "json": "json", "md": "md", "csv": "csv", "svg": "svg"}
    ext = ext_map.get(filetype, "txt")
    if not filename.endswith(f".{ext}"):
        filename = f"{filename}.{ext}"
    
    fpath = f"/app/uploads/{filename}"
    os.makedirs("/app/uploads", exist_ok=True)
    
    with open(fpath, "w") as f:
        f.write(content)
    
    return {
        "status": "ok",
        "filename": filename,
        "url": f"/download/{filename}",
        "size_bytes": len(content)
    }

@router.get("/api/tools/list-files")
async def list_files():
    """Список созданных файлов."""
    uploads = "/app/uploads"
    if not os.path.exists(uploads):
        return {"files": []}
    files = [{"name": f, "url": f"/download/{f}"} for f in os.listdir(uploads) if os.path.isfile(os.path.join(uploads, f))]
    return {"files": files, "count": len(files)}

@router.delete("/api/tools/delete-file/{filename}")
async def delete_file(filename: str):
    """Удаляет файл."""
    import os
    fpath = f"/app/uploads/{filename}"
    if os.path.exists(fpath):
        os.remove(fpath)
        return {"status": "deleted", "filename": filename}
    return {"error": "File not found"}

@router.get("/api/tools")
async def list_tools():
    """Возвращает список доступных инструментов для агентов."""
    from app.tools_schema import SKV_TOOLS
    return {"tools": SKV_TOOLS, "count": len(SKV_TOOLS)}
